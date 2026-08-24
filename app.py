# -*- coding: utf-8 -*-
"""
קטלוג חומרי הוראה למתמטיקה - גרסת ענן (רב-משתמשים)
======================================================
בשונה מהגרסה המקומית, האפליקציה הזו לא קוראת תיקייה מהמחשב -
היא מחוברת לחשבון Supabase (ענן חינמי) שמחזיק את כל הקבצים,
כדי שכל המורות/ים בצוות יוכלו להתחבר מכל מחשב, להעלות ולראות
את אותה ספרייה משותפת.

לפני הרצה יש לבצע הגדרה חד-פעמית לפי הקובץ SETUP_GUIDE.md.
"""

import os
import io
from datetime import datetime, timezone

import streamlit as st

try:
    from supabase import create_client, Client
except ImportError:
    create_client = None


# ---------------------------------------------------------------------------
# מילון הנושאים ומילות המפתח
# ---------------------------------------------------------------------------
TOPICS = {
    "שברים ואחוזים": [
        "שבר", "שברים", "מונה", "מכנה", "שבר מעורב", "אחוז", "אחוזים",
        "עשרוני", "עשרוניים", "צמצום", "הרחבה", "יחס", "פרופורציה", "חלק מתוך",
    ],
    "גאומטריה ומדידות": [
        "גאומטריה", "מצולע", "מצולעים", "משולש", "משולשים", "מרובע", "מרובעים",
        "ריבוע", "מלבן", "מקבילית", "מעוין", "טרפז", "מעגל", "זווית", "זוויות",
        "היקף", "שטח", "נפח", "סימטריה", "קודקוד", "צלע",
    ],
    "מספרים ופעולות": [
        "חיבור", "חיסור", "כפל", "חילוק", "לוח הכפל", "סדר פעולות", "סוגריים",
        "עשרות", "מאות", "אלפים", "זוגי", "אי זוגי", "חזקה", "שורש",
        "ערך המקום", "פתרון תרגיל", "משוואה", "תרגיל",
    ],
    "סטטיסטיקה ונתונים": [
        "דיאגרמה", "גרף", "טבלה", "ממוצע", "שכיח", "חציון", "הסתברות",
        "סקר", "שכיחות", "איסוף נתונים", "ייצוג נתונים",
    ],
    "חידות ומשחקים": [
        "חידה", "חידות", "משחק", "משחקים", "אתגר", "אולימפיאדה",
        "תשבץ", "סודוקו", "פאזל", "קואורדינטות", "חשיבה",
    ],
}
DEFAULT_TOPIC = "כללי - ללא סיווג"
ASSESSMENT_KEYWORDS = ["מבחן", "מבדק", "הערכה", "בוחן", "מבוא"]
ALL_TOPICS = list(TOPICS.keys()) + [DEFAULT_TOPIC]

MAX_CONTENT_CHARS = 6000
MAX_PDF_PAGES = 15
GEMINI_MODEL = "gemini-3.5-flash-lite"

BUCKET_NAME = "teaching-files"
TABLE_NAME = "files"

MIME_MAP = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".txt": "text/plain",
    ".csv": "text/csv",
    ".md": "text/markdown",
}


# ---------------------------------------------------------------------------
# עזרי הגדרות (secrets) - קוראים בזהירות כדי שלא יקרסו אם חסר משהו
# ---------------------------------------------------------------------------
def get_secret(key, default=None):
    try:
        return st.secrets.get(key, default)
    except Exception:
        return default


@st.cache_resource
def get_supabase_client():
    if create_client is None:
        return None
    url = get_secret("SUPABASE_URL")
    key = get_secret("SUPABASE_ANON_KEY")
    if not url or not key:
        return None
    try:
        return create_client(url, key)
    except Exception:
        return None


# ---------------------------------------------------------------------------
# חילוץ טקסט מבייטים של קובץ שהועלה (לצורך סיווג בלבד)
# ---------------------------------------------------------------------------
def extract_text_from_bytes(file_bytes, ext):
    ext = ext.lower()
    try:
        if ext in (".txt", ".md", ".csv"):
            return file_bytes.decode("utf-8", errors="ignore")[:MAX_CONTENT_CHARS]

        elif ext == ".docx":
            try:
                import docx
            except ImportError:
                return ""
            d = docx.Document(io.BytesIO(file_bytes))
            return "\n".join(p.text for p in d.paragraphs)[:MAX_CONTENT_CHARS]

        elif ext == ".pdf":
            try:
                from pypdf import PdfReader
            except ImportError:
                return ""
            reader = PdfReader(io.BytesIO(file_bytes))
            text = ""
            for i, page in enumerate(reader.pages):
                if i >= MAX_PDF_PAGES or len(text) >= MAX_CONTENT_CHARS:
                    break
                text += page.extract_text() or ""
            return text[:MAX_CONTENT_CHARS]

        elif ext == ".pptx":
            try:
                from pptx import Presentation
            except ImportError:
                return ""
            prs = Presentation(io.BytesIO(file_bytes))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                if len(text) >= MAX_CONTENT_CHARS:
                    break
            return text[:MAX_CONTENT_CHARS]

        elif ext in (".xlsx", ".xlsm"):
            try:
                import openpyxl
            except ImportError:
                return ""
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
            text = ""
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            text += str(cell) + " "
                    if len(text) >= MAX_CONTENT_CHARS:
                        break
                if len(text) >= MAX_CONTENT_CHARS:
                    break
            return text[:MAX_CONTENT_CHARS]

    except Exception:
        return ""
    return ""


# ---------------------------------------------------------------------------
# סיווג: שלב 1 מבוסס-חוקים, שלב 2 Gemini AI (רק כששם הקובץ כללי מדי)
# ---------------------------------------------------------------------------
def rule_based_topic(text):
    if not text:
        return None, []
    scores, matched = {}, {}
    for topic, keywords in TOPICS.items():
        hits = [kw for kw in keywords if kw in text]
        if hits:
            scores[topic] = len(hits)
            matched[topic] = hits
    if not scores:
        return None, []
    best = max(scores, key=scores.get)
    return best, matched[best]


def is_assessment(filename, content):
    combined = filename + " " + (content or "")
    return any(kw in combined for kw in ASSESSMENT_KEYWORDS)


def gemini_classify(filename, content, api_key):
    try:
        from google import genai
    except ImportError:
        return None, "הספרייה google-genai לא מותקנת"
    try:
        client = genai.Client(api_key=api_key)
        topics_list = "\n".join(f"- {t}" for t in ALL_TOPICS)
        prompt = (
            "אתה עוזר שמסווג חומרי הוראה במתמטיקה לפי נושא.\n"
            f"שם הקובץ: {filename}\n"
            f"קטע מתוך תוכן הקובץ (אם קיים):\n{(content or '')[:3000]}\n\n"
            "בחר נושא אחד בלבד מתוך הרשימה הבאה, והחזר אך ורק את שם הנושא "
            f"בדיוק כפי שהוא כתוב, ללא מילים נוספות:\n{topics_list}"
        )
        response = client.models.generate_content(model=GEMINI_MODEL, contents=prompt)
        answer = (getattr(response, "text", None) or "").strip()
        for topic in ALL_TOPICS:
            if topic in answer:
                return topic, None
        return None, f"תשובה לא ברורה מה-AI: {answer[:120]}"
    except Exception as e:
        return None, f"שגיאה בפנייה ל-Gemini: {e}"


def classify_file(name, content, use_ai, gemini_key):
    fname_topic, fname_matched = rule_based_topic(name)
    content_topic, content_matched = rule_based_topic(content) if content else (None, [])
    ai_note = None

    if fname_topic:
        topic, matched_kw, source = fname_topic, fname_matched, "מילת מפתח בשם הקובץ"
    elif use_ai and gemini_key:
        ai_topic, ai_note = gemini_classify(name, content, gemini_key)
        if ai_topic:
            topic, matched_kw, source = ai_topic, content_matched, "סיווג AI (Gemini)"
        else:
            topic = content_topic or DEFAULT_TOPIC
            matched_kw = content_matched
            source = "מילת מפתח בתוכן (AI נכשל)" if content_topic else "ללא התאמה (AI נכשל)"
    else:
        topic = content_topic or DEFAULT_TOPIC
        matched_kw = content_matched
        source = "מילת מפתח בתוכן" if content_topic else "ללא התאמה"

    doc_type = "מבחן / הערכה" if is_assessment(name, content) else "חומר תרגול"
    return topic, matched_kw, source, doc_type, ai_note


def build_excel_bytes(rows):
    import openpyxl
    from openpyxl.styles import Font, Alignment

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "קטלוג"
    ws.sheet_view.rightToLeft = True

    headers = ["שם קובץ", "נושא", "סוג", "מילות מפתח שזוהו", "מקור הסיווג", "הועלה על ידי"]
    ws.append(headers)
    for cell in ws[1]:
        cell.font = Font(bold=True)
        cell.alignment = Alignment(horizontal="right")

    for r in rows:
        ws.append([
            r.get("filename", ""), r.get("topic", ""), r.get("doc_type", ""),
            r.get("matched_keywords", ""), r.get("classification_source", ""),
            r.get("uploaded_by", ""),
        ])
        for cell in ws[ws.max_row]:
            cell.alignment = Alignment(horizontal="right")

    widths = {"A": 40, "B": 22, "C": 16, "D": 40, "E": 22, "F": 26}
    for col, w in widths.items():
        ws.column_dimensions[col].width = w

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# ממשק המשתמש
# ---------------------------------------------------------------------------
st.set_page_config(page_title="קטלוג חומרי הוראה למתמטיקה", page_icon="📚", layout="wide")
st.markdown(
    """
    <style>
    html, body, [class*="css"] { direction: rtl; }
    .stTextInput input, .stTextArea textarea { text-align: right; direction: rtl; }
    </style>
    """,
    unsafe_allow_html=True,
)

supabase = get_supabase_client()

if supabase is None:
    st.title("📚 קטלוג חומרי הוראה למתמטיקה")
    st.error(
        "החיבור לענן (Supabase) עדיין לא הוגדר על השרת הזה.\n\n"
        "יש להגדיר את SUPABASE_URL ו-SUPABASE_ANON_KEY ב-Secrets של Streamlit Cloud, "
        "לפי ההוראות בקובץ SETUP_GUIDE.md."
    )
    st.stop()

if "sb_session" not in st.session_state:
    st.session_state.sb_session = None
if "sb_user_email" not in st.session_state:
    st.session_state.sb_user_email = None


def do_login(email, password):
    try:
        res = supabase.auth.sign_in_with_password({"email": email, "password": password})
        st.session_state.sb_session = res.session
        st.session_state.sb_user_email = getattr(res.user, "email", email) if res.user else email
        return True, None
    except Exception as e:
        return False, str(e)


def do_signup(email, password):
    try:
        supabase.auth.sign_up({"email": email, "password": password})
        return True, None
    except Exception as e:
        return False, str(e)


def do_logout():
    try:
        supabase.auth.sign_out()
    except Exception:
        pass
    st.session_state.sb_session = None
    st.session_state.sb_user_email = None


# ---- מסך כניסה ----
if not st.session_state.sb_session:
    st.title("📚 קטלוג חומרי הוראה למתמטיקה")
    st.caption("כניסה לספרייה המשותפת של הצוות")

    mode = st.radio("בחירת פעולה", ["התחברות", "הרשמה"], horizontal=True, label_visibility="collapsed")
    email = st.text_input("אימייל")
    password = st.text_input("סיסמה", type="password")

    if mode == "התחברות":
        if st.button("התחברות", type="primary"):
            if not email or not password:
                st.warning("נא למלא אימייל וסיסמה")
            else:
                ok, err = do_login(email, password)
                if ok:
                    st.rerun()
                else:
                    st.error(f"ההתחברות נכשלה: {err}")
    else:
        if st.button("הרשמה", type="primary"):
            if not email or not password:
                st.warning("נא למלא אימייל וסיסמה")
            else:
                ok, err = do_signup(email, password)
                if ok:
                    st.success("נרשמת בהצלחה! (אם נדרש אימות אימייל - יש לאשר לפני התחברות) כעת אפשר לעבור ל'התחברות'.")
                else:
                    st.error(f"ההרשמה נכשלה: {err}")
    st.stop()


# ---- מסך ראשי (אחרי כניסה) ----
gemini_key = get_secret("GEMINI_API_KEY")

with st.sidebar:
    st.markdown(f"מחוברת/ר כ: **{st.session_state.sb_user_email}**")
    if st.button("התנתקות"):
        do_logout()
        st.rerun()
    st.divider()
    if not gemini_key:
        gemini_key = st.text_input("מפתח Gemini API (אופציונלי)", type="password")
    use_ai = st.checkbox("השתמש ב-AI לקבצים עם שם כללי", value=bool(gemini_key), disabled=not gemini_key)
    st.divider()
    with st.expander("📖 מילון הנושאים ומילות המפתח"):
        for topic, kws in TOPICS.items():
            st.markdown(f"**{topic}**")
            st.caption(", ".join(kws))
        st.markdown(f"**{DEFAULT_TOPIC}**")
        st.caption("קבצים שלא נמצאה עבורם אף התאמה")

st.title("📚 קטלוג חומרי הוראה למתמטיקה")
st.caption("הספרייה המשותפת של הצוות - נגישה מכל מחשב, לכל מי שמחובר/ת")

tab1, tab2 = st.tabs(["📥 העלאה וסיווג", "📂 הספרייה המשותפת"])

# ---- טאב 1: העלאה ----
with tab1:
    uploaded_files = st.file_uploader(
        "בחרו קובץ אחד או כמה קבצים להעלאה",
        accept_multiple_files=True,
        type=["pdf", "docx", "pptx", "xlsx", "txt", "csv", "md"],
    )

    if uploaded_files and st.button("🔍 סווג והעלה לספרייה", type="primary"):
        progress = st.progress(0.0)
        total = len(uploaded_files)
        ok_count, fail_notes = 0, []

        for i, uf in enumerate(uploaded_files):
            name = uf.name
            ext = os.path.splitext(name)[1]
            file_bytes = uf.getvalue()
            content = extract_text_from_bytes(file_bytes, ext)
            topic, matched_kw, source, doc_type, ai_note = classify_file(name, content, use_ai and bool(gemini_key), gemini_key)

            ts = datetime.now(timezone.utc).strftime("%Y%m%d%H%M%S")
            storage_path = f"{topic}/{ts}_{name}"

            try:
                supabase.storage.from_(BUCKET_NAME).upload(
                    path=storage_path,
                    file=file_bytes,
                    file_options={"content-type": MIME_MAP.get(ext.lower(), "application/octet-stream")},
                )
                supabase.table(TABLE_NAME).insert({
                    "filename": name,
                    "topic": topic,
                    "doc_type": doc_type,
                    "matched_keywords": ", ".join(matched_kw) if matched_kw else "",
                    "classification_source": source,
                    "storage_path": storage_path,
                    "uploaded_by": st.session_state.sb_user_email,
                }).execute()
                ok_count += 1
            except Exception as e:
                fail_notes.append(f"{name}: {e}")

            progress.progress((i + 1) / total)

        progress.empty()
        if ok_count:
            st.success(f"הועלו וסווגו בהצלחה {ok_count} מתוך {total} קבצים.")
        if fail_notes:
            st.warning("קבצים שלא עלו:\n" + "\n".join(fail_notes))

# ---- טאב 2: ספרייה משותפת ----
with tab2:
    topic_filter = st.selectbox("סינון לפי נושא", ["הכול"] + ALL_TOPICS)

    rows = []
    try:
        query = supabase.table(TABLE_NAME).select("*").order("uploaded_at", desc=True)
        if topic_filter != "הכול":
            query = query.eq("topic", topic_filter)
        res = query.execute()
        rows = res.data or []
    except Exception as e:
        st.error(f"שגיאה בטעינת הספרייה המשותפת: {e}")

    st.caption(f"נמצאו {len(rows)} קבצים")

    if rows:
        excel_bytes = None
        try:
            excel_bytes = build_excel_bytes(rows)
        except ImportError:
            pass
        if excel_bytes:
            st.download_button(
                "⬇️ הורדת הקטלוג כקובץ Excel",
                data=excel_bytes,
                file_name="קטלוג_נושאים.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )

    for row in rows:
        with st.container(border=True):
            c1, c2 = st.columns([5, 1])
            with c1:
                st.markdown(f"**{row.get('filename', '?')}**")
                st.caption(
                    f"נושא: {row.get('topic', '-')} | סוג: {row.get('doc_type', '-')} | "
                    f"הועלה ע\"י: {row.get('uploaded_by', '-')}"
                )
                if row.get("matched_keywords"):
                    st.caption(f"מילות מפתח שזוהו: {row.get('matched_keywords')}")
                if row.get("classification_source"):
                    st.caption(f"מקור הסיווג: {row.get('classification_source')}")
            with c2:
                storage_path = row.get("storage_path")
                if storage_path:
                    try:
                        signed = supabase.storage.from_(BUCKET_NAME).create_signed_url(storage_path, 3600)
                        url = None
                        if isinstance(signed, dict):
                            url = signed.get("signedURL") or signed.get("signedUrl") or signed.get("signed_url")
                        else:
                            url = getattr(signed, "signed_url", None)
                        if url:
                            st.link_button("פתיחה / הורדה", url, use_container_width=True)
                        else:
                            st.caption("אין קישור זמין")
                    except Exception:
                        st.caption("שגיאה ביצירת קישור")

    if not rows:
        st.info("עדיין אין קבצים בספרייה המשותפת. אפשר להתחיל להעלות בטאב הראשון.")
